import argparse
import os
import sys
import re
import time
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import translators as ts
from tqdm import tqdm

# Ensure standard output can handle utf-8 in case of console printing on Windows
sys.stdout.reconfigure(encoding='utf-8')


def get_runs_from_text_frame(text_frame, path):
    """Recursively extract text runs from a text frame."""
    runs = []
    for p_idx, paragraph in enumerate(text_frame.paragraphs):
        for r_idx, run in enumerate(paragraph.runs):
            if run.text.strip():
                runs.append({
                    'path': path + [p_idx, r_idx],
                    'run_obj': run,
                    'text': run.text
                })
    return runs


def iter_shapes(shapes, path=[]):
    """Recursively extract text runs from all shapes (including tables and groups)."""
    runs = []
    for s_idx, shape in enumerate(shapes):
        current_path = path + [s_idx]
        if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
            r = get_runs_from_text_frame(shape.text_frame, current_path)
            runs.extend(r)
        elif getattr(shape, 'has_table', False) and shape.has_table:
            for row_idx, row in enumerate(shape.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_path = current_path + ['table', row_idx, col_idx]
                    r = get_runs_from_text_frame(cell.text_frame, cell_path)
                    runs.extend(r)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            r = iter_shapes(shape.shapes, current_path + ['group'])
            runs.extend(r)
    return runs


def is_translatable(text):
    """Basic check to avoid translating pure numbers or symbols."""
    # If the text has letters or unicode characters (e.g., Chinese, Japanese, Arabic, etc.), we translate.
    # Exclude if it only contains digits, punctuations, and whitespace.
    return bool(re.search(r'[^\W\d_]', text, re.UNICODE))


def translate_presentation(input_path, output_path, from_lang, to_lang, translator_engine):
    print(f"Loading presentation: {input_path}")
    prs = Presentation(input_path)

    print("Extracting texts...")
    all_runs = []
    for slide_idx, slide in enumerate(prs.slides):
        runs = iter_shapes(slide.shapes, path=[slide_idx])
        for r in runs:
            if is_translatable(r['text']):
                all_runs.append(r)

    print(f"Found {len(all_runs)} text blocks to translate.")
    
    # We will track unique texts to minimize API calls
    unique_texts = {r['text'] for r in all_runs}
    translation_cache = {}

    print(f"Translating {len(unique_texts)} unique strings using '{translator_engine}'...")
    
    # Progress bar for unique texts
    for text in tqdm(unique_texts, desc="Translating", unit="str"):
        retries = 3
        while retries > 0:
            try:
                translated_text = ts.translate_text(
                    text, 
                    translator=translator_engine, 
                    from_language=from_lang, 
                    to_language=to_lang
                )
                translation_cache[text] = translated_text
                # Add a small delay to avoid hitting rate limits
                time.sleep(0.3)
                break
            except Exception as e:
                retries -= 1
                if retries == 0:
                    print(f"\n[Warning] Failed to translate: '{text}'. Error: {e}")
                    translation_cache[text] = text # Fallback to original text
                else:
                    time.sleep(1.0) # Wait before retry

    print("\nApplying translations to presentation...")
    for r in all_runs:
        original_text = r['text']
        if original_text in translation_cache:
            r['run_obj'].text = translation_cache[original_text]

    print(f"Saving translated presentation to: {output_path}")
    prs.save(output_path)
    print("Done!")


def main():
    parser = argparse.ArgumentParser(
        description="Translate PowerPoint (PPTX) files easily via command line.",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter
    )
    
    parser.add_argument("-i", "--input", required=True, help="Path to the input .pptx file")
    parser.add_argument("-o", "--output", help="Path to save the translated .pptx file. Defaults to <input>_translated.pptx")
    parser.add_argument("-f", "--from_lang", default="auto", help="Source language code (e.g., 'auto', 'zh', 'en', 'fr')")
    parser.add_argument("-t", "--to_lang", default="en", help="Target language code (e.g., 'en', 'zh', 'es', 'de')")
    parser.add_argument("-e", "--engine", default="bing", help="Translation engine to use (e.g., 'bing', 'google', 'alibaba')")

    args = parser.parse_args()

    input_file = args.input
    if not os.path.isfile(input_file):
        print(f"Error: The input file '{input_file}' does not exist.")
        sys.exit(1)

    if not input_file.lower().endswith('.pptx'):
        print(f"Error: The input file must be a .pptx file.")
        sys.exit(1)

    output_file = args.output
    if not output_file:
        base, ext = os.path.splitext(input_file)
        output_file = f"{base}_{args.to_lang}{ext}"

    try:
        translate_presentation(
            input_path=input_file,
            output_path=output_file,
            from_lang=args.from_lang,
            to_lang=args.to_lang,
            translator_engine=args.engine
        )
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
